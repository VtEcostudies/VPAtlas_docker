#!/usr/bin/env python3
"""Generate VPAtlas v3 presentation as .pptx for Google Slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette — VT Center for Ecostudies brand greens
DARK_BG = RGBColor(0x00, 0x55, 0x4A)        # VCE dark teal-green (primary brand)
ACCENT_GREEN = RGBColor(0x00, 0x74, 0x65)   # VCE medium teal-green
ACCENT_LIGHT = RGBColor(0x00, 0xD0, 0x84)   # vivid green accent
ACCENT_TEAL = RGBColor(0x00, 0xBF, 0xFF)    # cyan (pool status: probable)
ACCENT_GOLD = RGBColor(0xDA, 0xA5, 0x20)    # goldenrod (pool status: potential)
ACCENT_NAVY = RGBColor(0x00, 0x00, 0x8B)    # dark blue (pool status: confirmed)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
MED_GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
GREEN = RGBColor(0x00, 0x74, 0x65)           # VCE medium teal-green
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)
SLIDE_BG = RGBColor(0xF0, 0xF7, 0xF5)       # very light green tint


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=18,
                     color=DARK_TEXT, spacing=Pt(6), font_name="Calibri", sub_items=None,
                     bullet_color=None):
    """Add bulleted text with visual bullet characters.
    sub_items is a dict mapping item index to list of sub-bullets."""
    from pptx.oxml.ns import qn
    import copy

    if bullet_color is None:
        bullet_color = ACCENT_GREEN

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    def _set_bullet(paragraph, level=0, b_color=None):
        """Set XML-level bullet properties for a paragraph."""
        if b_color is None:
            b_color = bullet_color
        pPr = paragraph._pPr
        if pPr is None:
            pPr = paragraph._p.get_or_add_pPr()
        # Set indentation
        indent = Emu(228600) if level == 0 else Emu(457200)  # 0.25in / 0.5in
        margin = Emu(457200) if level == 0 else Emu(685800)  # 0.5in / 0.75in
        pPr.set('indent', str(-indent))
        pPr.set('marL', str(margin))
        # Add bullet character
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})  # bullet •
        if level == 1:
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2013'})  # en-dash –
        # Bullet color
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (b_color[0], b_color[1], b_color[2])})
        buClr.append(srgbClr)
        # Bullet size
        buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '100000'})
        # Remove any existing bullet settings
        for tag in [qn('a:buNone'), qn('a:buChar'), qn('a:buClr'), qn('a:buSzPct')]:
            for existing in pPr.findall(tag):
                pPr.remove(existing)
        pPr.append(buSzPct)
        pPr.append(buClr)
        pPr.append(buChar)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        _set_bullet(p, level=0)

        if sub_items and i in sub_items:
            for sub in sub_items[i]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.font.size = Pt(font_size - 2)
                sp.font.color.rgb = MED_GRAY
                sp.font.name = font_name
                sp.space_after = Pt(3)
                sp.level = 1
                _set_bullet(sp, level=1, b_color=MED_GRAY)

    return txBox


def add_section_header(slide, text, top=Inches(0.4)):
    add_text_box(slide, Inches(0.8), top, Inches(11), Inches(0.7),
                 text, font_size=32, color=DARK_BG, bold=True)


# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             "VPAtlas v3", font_size=60, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "Rewriting the UI/UX of the Vermont Vernal Pool Atlas",
             font_size=28, color=ACCENT_LIGHT)
add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.6),
             "Vermont Center for Ecostudies  |  April 2026",
             font_size=18, color=MED_GRAY)

# Five colored dots representing pool statuses
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)
for i, (clr, label) in enumerate([
    (ACCENT_GOLD, "Potential"), (ACCENT_TEAL, "Probable"), (ACCENT_NAVY, "Confirmed"),
    (GRAY, "Duplicate"), (DARK_RED, "Eliminated"),
]):
    x = Inches(1 + i * 1.8)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(5.5), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()
    add_text_box(slide, x + Inches(0.4), Inches(5.45), Inches(1.5), Inches(0.4),
                 label, font_size=14, color=LIGHT_GRAY)


# ─── SLIDE 2: What VPAtlas Does ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "What VPAtlas Does Today")

items = [
    "Authoritative live database of Vermont's vernal pools",
    "Tracks pools through the pipeline: Identified → Visited → Monitored",
    "Five pool statuses: Potential, Probable, Confirmed, Duplicate, Eliminated",
    "Ingests S123 field data via custom API process, manually invoked by Admins",
    "Publishes data via live API to VT ANR / VCGI database",
    "User self-registration and JWT-based login",
    "Users enter and edit Atlas Visit observations online",
    "Admins manage users, ingest S123 data, review visits, set pool statuses",
    "Integrates pool, visit, and survey data with photos (S123 visit photos broken)",
    "Public explore interface with interactive map",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                 items, font_size=20, color=DARK_TEXT)


# ─── SLIDE 3: Why v3 — Platform Problems ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "Why v3 — The UI Platform is Aging")

items = [
    "The DB and API are solid — it's the UI/UX that needs replacing",
    "UI built on Angular 14 — now 3+ major versions behind",
    "Heavy dependency chain: TypeScript, RxJS 6, Angular CLI, build tooling",
    "jQuery + DataTables mixed alongside Angular (anti-pattern)",
    "Deprecated dependencies: request, Protractor, TSLint, AWS SDK v2",
    "No containerization — PM2 process management, manual deployment",
    "Hard to update, hard to recruit help, fragile build process",
]
sub = {
    1: ["Angular ecosystem is moving toward signals/standalone — v14 patterns are legacy"],
    5: ["Each deployment is manual; no reproducible environment"],
}
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                 items, font_size=20, color=DARK_TEXT, sub_items=sub)


# ─── SLIDE 4: Why v3 — Data Collection ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "Why v3 — Data Collection is Broken")

# Two columns
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "In-App Visit Editing", font_size=22, color=DARK_BG, bold=True)
items_left = [
    "Desktop-oriented, not field-friendly",
    "Clunky multi-page forms",
    "No offline capability",
    "No GPS integration",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3),
                 items_left, font_size=18, color=DARK_TEXT)

add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
             "Survey123 Band-Aid", font_size=22, color=DARK_BG, bold=True)
items_right = [
    "Separate platform = separate login, data silo",
    "Form updates easily break the ingest pipeline",
    "No control over UX or behavior",
    "Visit photo integration is broken",
    "Better than nothing, but fragile",
]
add_bullet_slide(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(3),
                 items_right, font_size=18, color=DARK_TEXT)

# Bottom callout
add_shape_bg(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.8), DARK_BG)
add_text_box(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(0.6),
             "v3: Surveys write direct to DB. VCE controls the full stack. Future changes are easy.",
             font_size=20, color=ACCENT_LIGHT, bold=True)


# ─── SLIDE 5: Why v3 — The Opportunity ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "Why v3 — The Opportunity")

items = [
    "Proven tech from LoonWeb",
    "PWA: installable from browser, no app store, auto-updates",
    "Offline-first: IndexedDB data queue, smart sync on reconnect",
    "Offline map basemaps and overlays for field use without cell service",
    "Real-time GPS tracking with accuracy monitoring",
    "Full control over data-collection design (e.g. select a pool by name or on map)",
    "Add PoolFinder: GPS navigation to pools in the field",
    "Mobile-first responsive design with Bootstrap 5",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(3.5),
                 items, font_size=19, color=DARK_TEXT)

# Urgency callout
add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(1.2), Inches(5.1), Inches(11), Inches(0.5),
             "Wetlands Act Changes", font_size=22, color=ACCENT_GOLD, bold=True)
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
             "With imminent regulatory changes, we need tools for rapid identification of potential pools",
             font_size=18, color=WHITE)


# ─── SLIDE 6: Architecture Comparison ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "Architecture: v2 → v3")

# Table
rows_data = [
    ("", "VPAtlas v2", "VPAtlas v3"),
    ("UI Framework", "Angular 14, TypeScript, build chain", "Plain HTML/JS/CSS, ES6 modules, no build step"),
    ("Deployment", "PM2, manual server setup", "Docker Compose (db + api + ui)"),
    ("Database", "PostgreSQL + PostGIS (reused)", "PostgreSQL 17 + PostGIS 3.5 (migrated, same schema)"),
    ("API", "Express/Node (reused)", "Same API, enhanced with env-var config overlay"),
    ("State Mgmt", "Angular services, RxJS", "URL params + IndexedDB persistence"),
    ("Maps", "Leaflet + Esri", "Leaflet + VCGI tiles + boundary overlays"),
    ("Data Collection", "S123 + clunky web forms", "Integrated PWA with GPS + offline"),
]

table_shape = slide.shapes.add_table(len(rows_data), 3,
                                      Inches(0.8), Inches(1.5),
                                      Inches(11.7), Inches(5.0))
table = table_shape.table

# Column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.85)
table.columns[2].width = Inches(4.85)

for row_idx, row_data in enumerate(rows_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.name = "Calibri"
            if row_idx == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK_TEXT
            if col_idx == 0:
                paragraph.font.bold = True

        # Header row styling
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)


# ─── SLIDE 7: What's Built — Explore ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "What's Built — Explore Interface")

items = [
    "Three-pane layout: pool list | map | summary",
    "Filter system with data-type buttons, pool ID type-ahead, town/county multi-select tokens, status checkboxes",
    "All filters persist to IndexedDB, restore on page load, sync to URL",
    "Map markers: color = status (goldenrod / cyan / dark blue), shape = survey level (circle / triangle / diamond)",
    "VCGI basemaps: Color Infrared, Leaf-Off, Lidar DEM/DSM/Slope",
    "Clickable county and town boundary overlays — click to zoom and filter",
    "Pool detail view with visit and survey history",
    "61 automated tests via test_stack.sh",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                 items, font_size=20, color=DARK_TEXT)



# ─── SLIDE 8: What's Built — Admin & Auth ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "What's Built — Admin & Auth")

items = [
    "JWT-based user registration and login",
    "Review creation form with pool selection",
    "Review list and detail views",
    "User administration panel",
    "Profile management",
    "Dockerized stack: one command to start everything",
    "Database restore from backup with db_restore.sh",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                 items, font_size=20, color=DARK_TEXT)


# ─── SLIDE 9: What's Coming — Field App ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "What's Coming — Integrated Field App")

# Three feature columns
features = [
    ("PoolFinder", ACCENT_GREEN, [
        "GPS navigation to pools",
        "Live position on map",
        "Navigate to nearest pool",
        "Launch Visit or Survey",
        "The /survey landing page",
    ]),
    ("Visit & Survey Forms", DARK_BG, [
        "Mobile-first multi-page forms",
        "5-page Visit form",
        "Monitoring survey from S123 spec",
        "Photo capture & attachment",
        "Offline data queue + sync",
    ]),
    ("GPS Tools", ACCENT_GOLD, [
        "Pool boundary mapping",
        "Walk perimeter → polygon",
        "In-app polygon draw tools",
        "Accuracy monitoring",
        "Track recording",
    ]),
]

for i, (title, color, bullets) in enumerate(features):
    x = Inches(0.8 + i * 4.1)
    add_shape_bg(slide, x, Inches(1.5), Inches(3.7), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.2), Inches(1.5), Inches(3.3), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True)
    add_bullet_slide(slide, x + Inches(0.2), Inches(2.2), Inches(3.3), Inches(4),
                     bullets, font_size=16, color=DARK_TEXT)


# ─── SLIDE 10: PWA & Mobile Details ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "PWA & Mobile Architecture")

items = [
    "Progressive Web App — installable from browser, no app store required",
    "Service worker with multi-tier caching: app assets, data, map tiles",
    "Offline basemaps and boundary overlays for field use",
    "App versioning with build-time injection (sw_template.js pattern from LoonWeb)",
    "BroadcastChannel communication between service worker and app",
    "IndexedDB for all local state: filters, survey data, sync queue",
    "iOS-compatible: silent audio keep-alive for background GPS",
    "Capacitor bridge available for App Store distribution if needed",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                 items, font_size=20, color=DARK_TEXT)


# ─── SLIDE 11: Demo ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Live Demo", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "VPAtlas v3 Explore Interface", font_size=24, color=ACCENT_LIGHT)

demo_items = [
    "Three-pane pool explorer with filters",
    "Map with VCGI basemaps and boundary overlays",
    "Pool detail with visit/survey history",
    "Mobile responsive layout",
]
add_bullet_slide(slide, Inches(1), Inches(4.8), Inches(11), Inches(2.5),
                 demo_items, font_size=18, color=LIGHT_GRAY)


# ─── SLIDE 12: Summary ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_section_header(slide, "Summary")

items = [
    "v2 works — the DB and API are solid, but the UI/UX is built on aging tech",
    "v3 reuses DB and API, replaces the Angular UI with plain modern JS",
    "Simpler, maintainable, no build chain — Dockerized deployment",
    "Integrated data-collection: surveys direct to DB, VCE controls the full stack",
    "Mobile / offline / GPS capabilities proven across VCE projects",
    "Positioned for wetlands act changes with rapid pool identification tools",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(4),
                 items, font_size=22, color=DARK_TEXT, spacing=Pt(12))

# Bottom accent
add_shape_bg(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), DARK_BG)
add_text_box(slide, Inches(1), Inches(7.05), Inches(11), Inches(0.4),
             "Vermont Center for Ecostudies  •  VPAtlas v3  •  April 2026",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# Save
output_path = "/home/jloomis/VPAtlas/VPAtlas_docker/presentations/VPAtlas_v3_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
